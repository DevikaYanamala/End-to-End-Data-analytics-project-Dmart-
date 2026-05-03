"""
DMart Opportunities Analysis - PowerPoint Presentation Builder
Creates a 10-slide presentation based on Food Services brief
"""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

BASE = os.path.dirname(__file__)
CHARTS = os.path.join(BASE, 'charts')
DATA_FILE = os.path.join(BASE, 'DMart Data Store.xlsx')

# Load data for stats
df = pd.read_excel(DATA_FILE, sheet_name='Order Data')
df['Year'] = df['Order Date'].dt.year

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_text(slide, left, top, width, height, bullets, size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return tf

def add_chart_img(slide, chart_name, left=0.5, top=1.5, width=12.3, height=5.5):
    path = os.path.join(CHARTS, chart_name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

def slide_title_bar(slide, title, subtitle=None):
    add_textbox(slide, 0.5, 0.3, 12, 0.7, title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.5, 0.95, 12, 0.5, subtitle, size=14, color=RGBColor(0xBB,0xBB,0xBB))

# Compute key stats
total_sales = df['Sales'].sum()
total_profit = df['Profit'].sum()
total_orders = df['Order ID'].nunique()
total_customers = df['Customer Name'].nunique()
avg_margin = (total_profit / total_sales * 100)
yoy_14 = ((df[df['Year']==2014]['Sales'].sum() / df[df['Year']==2013]['Sales'].sum()) - 1) * 100

# ============ SLIDE 1: Title Slide ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)
add_textbox(slide1, 1, 1.5, 11, 1.2, 'DMart Opportunities Analysis', size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slide1, 1, 2.8, 11, 0.8, 'Food Services Division — Strategic Growth Review', size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide1, 1, 4.0, 11, 0.5, 'Prepared for: Shaun, Sales Director', size=16, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
add_textbox(slide1, 1, 4.6, 11, 0.5, f'Data Period: 2011–2014  |  {total_orders:,} Orders  |  {total_customers} Customers  |  15 Countries', size=14, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
# divider line
shape = slide1.shapes.add_shape(1, Inches(3), Inches(3.8), Inches(7.33), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
print("Slide 1: Title - done")

# ============ SLIDE 2: Executive Summary ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
slide_title_bar(slide2, 'Executive Summary', 'Key findings at a glance')

metrics = [
    (f'€{total_sales/1e6:.1f}M', 'Total Sales'),
    (f'€{total_profit/1e3:.0f}K', 'Total Profit'),
    (f'{avg_margin:.1f}%', 'Profit Margin'),
    (f'+{yoy_14:.1f}%', 'YoY Growth (2014)'),
]
for i, (val, label) in enumerate(metrics):
    x = 0.8 + i * 3.1
    box = slide2.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(2.7), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    box.line.fill.background()
    add_textbox(slide2, x+0.1, 1.9, 2.5, 0.7, val, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide2, x+0.1, 2.5, 2.5, 0.4, label, size=13, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)

bullets = [
    "▸ Consumer segment dominates at 51.8% of sales, with Corporate (31.2%) as the biggest growth opportunity",
    "▸ Central region accounts for 55% of revenue — North and South regions are under-penetrated",
    "▸ France, Germany, and UK are top 3 markets — 7 smaller countries represent untapped potential",
    "▸ Discounts above 20% destroy profitability — disciplined pricing is critical for margin protection",
    "▸ Technology and Furniture categories offer higher margins vs. Office Supplies volume play",
    "▸ Year-over-year growth accelerating: from 12.7% (2012) to 25.3% (2014), signalling strong momentum"
]
add_bullet_text(slide2, 0.8, 3.5, 11.5, 3.5, bullets, size=13, color=WHITE)
print("Slide 2: Executive Summary - done")

# ============ SLIDE 3: Segment Performance ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, LIGHT_GRAY)
slide_title_bar(slide3, 'Segment Performance Overview', 'Share of total spends across Customer, Corporate & Home Office')
for t in slide3.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide3, 'chart1_segment_overview.png', 0.5, 1.6, 12.3, 5.5)
print("Slide 3: Segment Performance - done")

# ============ SLIDE 4: Category Spend by Segment ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, LIGHT_GRAY)
slide_title_bar(slide4, 'Category Spend Share by Segment', 'How each segment allocates spend across product categories')
for t in slide4.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide4, 'chart2_category_by_segment.png', 0.5, 1.6, 12.3, 5.5)
print("Slide 4: Category Spend - done")

# ============ SLIDE 5: Geographic Distribution ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, LIGHT_GRAY)
slide_title_bar(slide5, 'Geographic Distribution of Sales', 'Sales breakdown across 15 European countries and 3 regions')
for t in slide5.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide5, 'chart3_geographic_sales.png', 0.5, 1.6, 12.3, 5.5)
print("Slide 5: Geographic - done")

# ============ SLIDE 6: Growth Trends ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, LIGHT_GRAY)
slide_title_bar(slide6, 'Sales Growth Trend (2011–2014)', 'Year-over-year growth acceleration signals strong momentum')
for t in slide6.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide6, 'chart4_yoy_growth.png', 0.5, 1.6, 12.3, 5.5)
print("Slide 6: Growth Trends - done")

# ============ SLIDE 7: Opportunity Matrix ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, LIGHT_GRAY)
slide_title_bar(slide7, 'Sub-Category Opportunity Matrix', 'Identifying high-margin, high-volume growth opportunities')
for t in slide7.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide7, 'chart5_opportunity_matrix.png', 0.5, 1.5, 12.3, 5.6)
print("Slide 7: Opportunity Matrix - done")

# ============ SLIDE 8: Top Cities + Discount Impact ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, LIGHT_GRAY)
slide_title_bar(slide8, 'City Performance & Discount Impact', 'Where we win and how discounting affects profitability')
for t in slide8.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide8, 'chart6_top_cities.png', 0.2, 1.5, 6.3, 5.0)
add_chart_img(slide8, 'chart7_discount_impact.png', 6.6, 1.5, 6.5, 5.0)
print("Slide 8: Cities & Discounts - done")

# ============ SLIDE 9: Adjacent Opportunities ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, LIGHT_GRAY)
slide_title_bar(slide9, 'Adjacent Segment & Regional Opportunities', 'Identifying under-penetrated segments across geographies')
for t in slide9.shapes:
    if t.has_text_frame:
        for p in t.text_frame.paragraphs:
            p.font.color.rgb = DARK_TEXT
add_chart_img(slide9, 'chart8_region_segment_heatmap.png', 0.2, 1.5, 6.0, 5.0)
add_chart_img(slide9, 'chart9_penetration_opportunity.png', 6.4, 1.5, 6.7, 5.0)
print("Slide 9: Adjacent Opportunities - done")

# ============ SLIDE 10: Recommendations & Next Steps ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)
slide_title_bar(slide10, 'Strategic Recommendations & Next Steps', 'Where are our biggest opportunities?')

recs = [
    "1. EXPAND CORPORATE SEGMENT — Corporate represents 31% of sales but shows higher avg. order values. "
    "Targeted campaigns in France, Germany & UK can capture an additional estimated €200K+ annually.",
    "",
    "2. PENETRATE UNDER-SERVED GEOGRAPHIES — 7 smaller markets (Nordics, Benelux, Portugal, Switzerland) "
    "account for only 14% of sales. Deploying dedicated sales coverage can unlock significant new revenue.",
    "",
    "3. OPTIMIZE DISCOUNT STRATEGY — Discounts above 20% result in negative profit margins. Implement "
    "tiered discount caps: max 10% for Office Supplies, max 15% for Technology, max 20% for Furniture.",
    "",
    "4. DOUBLE DOWN ON HIGH-MARGIN SUB-CATEGORIES — Copiers, Phones, and Accessories show the best "
    "margin-to-volume ratio. Prioritize these in sales incentive plans.",
    "",
    "5. INVEST IN NORTH & SOUTH REGIONS — Central region is mature at 55% share. North (22%) and "
    "South (23%) offer the best incremental growth with focused regional sales teams.",
    "",
    "6. LEVERAGE CUSTOMER FEEDBACK — Only 50% of transactions include feedback. Implement systematic "
    "feedback collection to improve NPS and retention in top accounts."
]
add_bullet_text(slide10, 0.8, 1.8, 11.5, 5.5, recs, size=13, color=WHITE)
print("Slide 10: Recommendations - done")

# Save
output_path = os.path.join(BASE, 'DMart_Opportunities_Analysis.pptx')
prs.save(output_path)
print(f"\nPresentation saved: {output_path}")
